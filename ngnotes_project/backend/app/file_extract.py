"""
NGNotes: File extraction utilities for raw engineering notes.
"""

from __future__ import annotations

import csv
import json
import os
import subprocess
import tempfile
import xml.etree.ElementTree as ET
from html.parser import HTMLParser
from pathlib import Path
from typing import Tuple

from fastapi import UploadFile
from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {
    # Documents
    ".pdf", ".doc", ".docx", ".pptx", ".rtf",
    # Spreadsheets
    ".xlsx", ".csv", ".tsv",
    # Plain text and structured text
    ".md", ".txt", ".json", ".xml", ".log",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    # Web
    ".html", ".htm",
    # Diagram / "code-as-diagram" formats (text under the hood)
    ".svg", ".mmd", ".mermaid", ".puml", ".plantuml", ".dot", ".gv", ".drawio",
}


def _normalize_text(text: str) -> str:
    lines = [line.rstrip() for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")]
    # Keep paragraph boundaries while preventing very long blank runs.
    out = []
    blank_count = 0
    for line in lines:
        if line.strip() == "":
            blank_count += 1
            if blank_count <= 1:
                out.append("")
        else:
            blank_count = 0
            out.append(line)
    return "\n".join(out).strip()


def _extract_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    chunks = []
    for page in reader.pages:
        chunks.append(page.extract_text() or "")
    return _normalize_text("\n\n".join(chunks))


def _extract_json(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
        obj = json.load(fh)
    return _normalize_text(json.dumps(obj, indent=2, ensure_ascii=False))


def _extract_xml(file_path: str) -> str:
    tree = ET.parse(file_path)
    root = tree.getroot()
    text = ET.tostring(root, encoding="unicode", method="text")
    return _normalize_text(text)


def _extract_svg(file_path: str) -> str:
    """SVG is XML. Pull readable text plus title/desc elements if present."""
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
    except ET.ParseError:
        # Some SVG exports include odd byte order marks or comments; fall back to plain read.
        return _extract_text_plain(file_path)

    pieces: list[str] = []
    # Collect title/desc up front so the model gets the diagram's stated intent.
    for tag in ("title", "desc"):
        for el in root.iter():
            if el.tag.endswith(tag) and (el.text or "").strip():
                pieces.append(f"[{tag}] {el.text.strip()}")

    body_text = ET.tostring(root, encoding="unicode", method="text")
    if body_text and body_text.strip():
        pieces.append(body_text)

    return _normalize_text("\n".join(pieces))


def _extract_pptx(file_path: str) -> str:
    """Extract slide text from a .pptx using python-pptx."""
    from pptx import Presentation  # local import keeps module load cheap

    prs = Presentation(file_path)
    chunks: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = [f"## Slide {idx}"]

        # Slide title (if any)
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None and getattr(title_shape, "has_text_frame", False):
            t = (title_shape.text_frame.text or "").strip()
            if t:
                slide_lines.append(f"# {t}")

        for shape in slide.shapes:
            if shape == title_shape:
                continue
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                txt = (shape.text_frame.text or "").strip()
                if txt:
                    slide_lines.append(txt)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [(c.text or "").strip() for c in row.cells]
                    if any(cells):
                        slide_lines.append(" | ".join(cells))

        notes = getattr(slide, "notes_slide", None)
        if notes is not None:
            note_text = (notes.notes_text_frame.text or "").strip() if notes.notes_text_frame else ""
            if note_text:
                slide_lines.append(f"_Speaker notes: {note_text}_")

        chunks.append("\n".join(slide_lines))

    return _normalize_text("\n\n".join(chunks))


def _extract_text_plain(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
        return _normalize_text(fh.read())


def _extract_doc_or_docx_via_textutil(file_path: str) -> str:
    # macOS textutil handles .doc, .docx, .rtf without extra runtime deps.
    proc = subprocess.run(
        ["textutil", "-convert", "txt", "-stdout", file_path],
        capture_output=True,
        text=True,
        check=False,
    )
    if proc.returncode != 0:
        raise RuntimeError(proc.stderr.strip() or "Failed to convert document via textutil")
    return _normalize_text(proc.stdout)


def _extract_xlsx(file_path: str) -> str:
    """Render each worksheet of an .xlsx as '## Sheet: <name>' + ' | '-joined rows."""
    from openpyxl import load_workbook  # local import keeps module load cheap

    wb = load_workbook(file_path, data_only=True, read_only=True)
    chunks: list[str] = []
    for ws in wb.worksheets:
        sheet_lines: list[str] = [f"## Sheet: {ws.title}"]
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v).strip() for v in row]
            if not any(cells):
                continue
            sheet_lines.append(" | ".join(cells))
        if len(sheet_lines) > 1:
            chunks.append("\n".join(sheet_lines))
    wb.close()
    return _normalize_text("\n\n".join(chunks))


def _extract_delimited(file_path: str, delimiter: str) -> str:
    """Read a CSV/TSV-style file into ' | '-joined rows."""
    with open(file_path, "r", encoding="utf-8", errors="ignore", newline="") as fh:
        reader = csv.reader(fh, delimiter=delimiter)
        rows = [" | ".join(cell.strip() for cell in row) for row in reader if any(c.strip() for c in row)]
    return _normalize_text("\n".join(rows))


class _HTMLTextExtractor(HTMLParser):
    """Collect visible text from an HTML document, skipping script/style and
    inserting blank lines around block-level elements."""

    _SKIP_TAGS = {"script", "style", "noscript"}
    _BLOCK_TAGS = {
        "p", "div", "section", "article", "header", "footer", "main", "aside",
        "h1", "h2", "h3", "h4", "h5", "h6",
        "ul", "ol", "li", "table", "tr", "td", "th",
        "br", "hr", "pre", "blockquote", "figure",
    }

    def __init__(self) -> None:
        super().__init__()
        self._buf: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs) -> None:
        tag = tag.lower()
        if tag in self._SKIP_TAGS:
            self._skip_depth += 1
        elif tag in self._BLOCK_TAGS:
            self._buf.append("\n")

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag in self._SKIP_TAGS and self._skip_depth > 0:
            self._skip_depth -= 1
        elif tag in self._BLOCK_TAGS:
            self._buf.append("\n")

    def handle_data(self, data: str) -> None:
        if self._skip_depth:
            return
        if data:
            self._buf.append(data)

    def text(self) -> str:
        return "".join(self._buf)


def _extract_html(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
        parser = _HTMLTextExtractor()
        parser.feed(fh.read())
        parser.close()
    return _normalize_text(parser.text())


def extract_text_from_upload(upload: UploadFile) -> Tuple[str, str]:
    """Extract text from supported uploaded files.

    Returns (filename, extracted_text).
    """
    filename = upload.filename or "uploaded_file"
    suffix = Path(filename).suffix.lower()

    if suffix not in SUPPORTED_EXTENSIONS:
        pretty = sorted(s.lstrip(".") for s in SUPPORTED_EXTENSIONS)
        raise ValueError(
            f"Unsupported file type '{suffix or 'unknown'}'. Supported: {', '.join(pretty)}"
        )

    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        data = upload.file.read()
        tmp.write(data)
        tmp_path = tmp.name

    try:
        if suffix == ".pdf":
            text = _extract_pdf(tmp_path)
        elif suffix in {
            ".md", ".txt", ".log",
            ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
            ".mmd", ".mermaid", ".puml", ".plantuml", ".dot", ".gv",
        }:
            text = _extract_text_plain(tmp_path)
        elif suffix in {".doc", ".docx", ".rtf"}:
            text = _extract_doc_or_docx_via_textutil(tmp_path)
        elif suffix == ".pptx":
            text = _extract_pptx(tmp_path)
        elif suffix == ".xlsx":
            text = _extract_xlsx(tmp_path)
        elif suffix == ".csv":
            text = _extract_delimited(tmp_path, ",")
        elif suffix == ".tsv":
            text = _extract_delimited(tmp_path, "\t")
        elif suffix == ".json":
            text = _extract_json(tmp_path)
        elif suffix == ".svg":
            text = _extract_svg(tmp_path)
        elif suffix in {".html", ".htm"}:
            text = _extract_html(tmp_path)
        elif suffix in {".xml", ".drawio"}:
            # .drawio is XML with embedded mxGraph nodes; raw XML text extraction
            # still gives the model the labels and structure to reason about.
            text = _extract_xml(tmp_path)
        else:
            text = _extract_text_plain(tmp_path)

        if not text.strip():
            raise ValueError("No readable text was found in the uploaded file")

        return filename, text
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass
